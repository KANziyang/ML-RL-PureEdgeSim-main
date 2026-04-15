#!/usr/bin/env python3
"""
Generate comprehensive PowerPoint presentation from EE5003 Report
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# Initialize presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle=""):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = title
    if subtitle:
        subtitle_shape.text = subtitle

    return slide

def add_content_slide(prs, title, content_items):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]

    title_shape.text = title
    tf = body_shape.text_frame
    tf.clear()

    for item in content_items:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(18)

    return slide

def add_image_slide(prs, title, image_path, caption=""):
    """Add a slide with an image"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_shape = slide.shapes.title
    title_shape.text = title

    if os.path.exists(image_path):
        left = Inches(1)
        top = Inches(1.5)
        height = Inches(5)
        pic = slide.shapes.add_picture(image_path, left, top, height=height)

        if caption:
            left = Inches(1)
            top = Inches(6.8)
            width = Inches(8)
            height = Inches(0.5)
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.text = caption
            text_frame.paragraphs[0].font.size = Pt(12)

    return slide

def add_two_column_slide(prs, title, left_content, right_content):
    """Add a two-column content slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_shape = slide.shapes.title
    title_shape.text = title

    # Left column
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(4.5)
    height = Inches(5.5)
    textbox_left = slide.shapes.add_textbox(left, top, width, height)
    tf_left = textbox_left.text_frame
    tf_left.word_wrap = True
    for item in left_content:
        p = tf_left.add_paragraph()
        p.text = item
        p.font.size = Pt(14)

    # Right column
    left = Inches(5.2)
    textbox_right = slide.shapes.add_textbox(left, top, width, height)
    tf_right = textbox_right.text_frame
    tf_right.word_wrap = True
    for item in right_content:
        p = tf_right.add_paragraph()
        p.text = item
        p.font.size = Pt(14)

    return slide

# Base path for figures
fig_base = "EE5003_Report/figures/"

print("Creating presentation...")

# Slide 1: Title
add_title_slide(prs,
    "ON POLICY MULTI-AGENT REINFORCEMENT LEARNING\nFOR 6G EDGE-CLOUD COMPUTING",
    "Kan Ziyang\nA0326643A\nSupervisor: Prof. Tham Chen Khong\nMarch 2026")

# Slide 2: Abstract
add_content_slide(prs, "Abstract", [
    "Task offloading in 6G edge-cloud systems must coordinate execution placement and wireless PRB reservation",
    "Proposed MAPPO scheduler for joint task offloading and PRB allocation in LOCAL-EDGE-CLOUD architecture",
    "MAPPO achieves 99.96% success rate vs 90.39% for best heuristic (ROUND_ROBIN)",
    "PPO baseline reaches 99.94%, showing on-policy RL effectiveness",
    "Adaptive PRB control is the dominant source of performance gain",
    "MAPPO provides narrower reliability-oriented advantage over PPO"
])

# Slide 3: Research Questions
add_content_slide(prs, "Research Questions", [
    "1. Can on-policy RL improve QoS and system efficiency over heuristic schedulers when destination and PRB are jointly optimized?",
    "2. Is adaptive PRB allocation the dominant source of performance gain relative to other MAPPO design choices?",
    "3. Does MAPPO provide meaningful advantage over shared-policy PPO baseline in deadline-sensitive reliability?"
])

# Slide 4: Main Contributions
add_content_slide(prs, "Main Contributions", [
    "Joint task-offloading and PRB-allocation formulated as Dec-POMDP",
    "MAPPO framework with CTDE: per-device actors + centralized critic",
    "Shared-policy PPO baseline under same environment interface",
    "Extended PureEdgeSim with PRB-aware networking and Java-Python co-simulation",
    "Comprehensive evaluation against heuristic and learning baselines"
])

# Slide 5: Background - 6G and MEC
add_content_slide(prs, "Background: 6G and Mobile Edge Computing", [
    "6G integrates intelligence, communication, and distributed computing",
    "MEC reduces latency by moving computation closer to users",
    "Radio-resource allocation is central to edge-cloud orchestration",
    "Scheduler must jointly determine WHERE and HOW MUCH wireless resource",
    "Without reasonable PRB allocation, attractive compute destination can produce poor outcome"
])

# Slide 6: Why Multi-Agent RL
add_content_slide(prs, "Why Multi-Agent Reinforcement Learning?", [
    "RL improves through interaction rather than analytical model",
    "Devices interact through shared PRBs, edge queues, and compute resources",
    "One device's choice affects others - naturally multi-agent problem",
    "MARL enables local decision-making while learning cooperative behavior",
    "Partial observability + shared constraints + system trade-offs"
])

# Slide 7: Related Work Comparison Table
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Related Work Comparison"
add_image_slide(prs, "Related Work Comparison",
    fig_base + "related_work_table.png",
    "This work: MAPPO/PPO with Local+Edge+Cloud, per-device agents, PRB action, centralized critic, 160 devices")

# Slide 8: System Architecture
add_image_slide(prs, "Three-Tier System Architecture",
    fig_base + "network_topology.png",
    "200m × 200m area with 160 devices, 4 edge data centers, 1 cloud node")

# Slide 9: PureEdgeSim Extensions
add_image_slide(prs, "Extended PureEdgeSim Architecture",
    fig_base + "extended_pureedgesim_rl_architecture.png",
    "Four coordinated layers: simulator extensions, RL interface, runtime support, telemetry")

# Slide 10: Task Orchestration Flow
add_image_slide(prs, "Event-Driven Task Lifecycle",
    fig_base + "task_orchestration_flow.png",
    "Decisions triggered when tasks arrive, preserving event-driven semantics")

# Slide 11: Java-Python Co-Simulation
add_image_slide(prs, "Java-Python Training Architecture",
    fig_base + "java_python_training_architecture.png",
    "Training mode: Python launches Java, synchronous TCP JSON protocol")

# Slide 12: Simulation Parameters
add_two_column_slide(prs, "Simulation Configuration", [
    "• Area: 200m × 200m",
    "• Devices: 160 (30% smartphone, 30% sensor)",
    "• Simulation time: 30 minutes",
    "• Mobile speed: 1.4 m/s",
    "• Edge coverage: 120m",
    "• Cloud distance: 160m"
], [
    "• WLAN bandwidth: 5000 Mbps",
    "• Total PRB blocks: 5000",
    "• Max PRBs per task: 50",
    "• Reference distance d₀: 20m",
    "• Distance exponent α: 0.5",
    "• VM policy: SPACE_SHARED"
])

# Slide 13: Communication Model
add_content_slide(prs, "PRB-Based Communication Model", [
    "Total capacity: 5000 Mbps discretized into 5000 PRB blocks",
    "Bandwidth = (BW_WLAN / PRB_total) × blocks × min(1, (d₀/max(d,d₀))^α)",
    "RL algorithms: fixed PRB reservation before transmission",
    "8 discrete PRB bins: {0.02, 0.05, 0.10, 0.20, 0.40, 0.60, 0.80, 1.00}",
    "Non-RL algorithms: dynamic PRB allocation by simulator",
    "Cloud uses fixed equivalent distance of 160m"
])

# Slide 14: Energy Model
add_content_slide(prs, "Energy Model", [
    "Wireless energy: EEZC protocol with free-space and multipath models",
    "E_tx = E_elec·B + E_fs·B·d² (d ≤ d₀) or E_mp·B·d⁴ (d > d₀)",
    "CPU energy: linear utilization model",
    "ΔE_cpu = [P_idle + (P_max - P_idle)·u] / 3600 · Δt",
    "Online standardization via EMA for RL reward stability",
    "Total energy = CPU energy + wireless energy"
])

# Slide 15: Dec-POMDP Formulation
add_content_slide(prs, "Dec-POMDP Formulation", [
    "Agent set N: all task-generating devices",
    "Observation o_i: task demand, device state, destination features, feasibility mask",
    "Action a_i = (dest, prb): execution destination + PRB reservation level",
    "Global state s: aggregated load, energy, PRB usage",
    "Objective: minimize weighted cost of failures, latency, energy, PRB, fallback",
    "Constraints: feasible destinations, PRB budget, deadline satisfaction"
])

# Slide 16: MAPPO Framework
add_image_slide(prs, "MAPPO CTDE Framework",
    fig_base + "mappo_ctde_framework.png",
    "Centralized training with decentralized execution")

# Slide 17: Observation Design
add_content_slide(prs, "Observation and State Design", [
    "Agent observation (13D): task demand (5) + device state (4) + network context (4)",
    "Destination features (11D per dest): load, capacity, delay, distance, tier indicators",
    "Destination mask: suppresses infeasible destinations",
    "PPO critic state (27D): agent obs + aggregated system statistics",
    "MAPPO critic state (41D): base state + destination distribution + PRB distribution",
    "All features normalized to [0,1] or [-2,2]"
])

# Slide 18: Action Space
add_content_slide(prs, "Action Space and Fallback", [
    "Composite action: (destination, PRB bin)",
    "Destination: local(0) → edge DCs (sorted by ID) → cloud",
    "PRB bins: 8 discrete levels mapping to {1, 2, 5, 10, 20, 30, 40, 50} blocks",
    "Local execution forces PRB=0, log-prob and entropy zeroed",
    "Destination fallback: if infeasible at execution, use shortest estimated completion",
    "Fallback transitions excluded from PPO updates, penalized in reward"
])

# Slide 19: Reward Function
add_content_slide(prs, "Reward Function Design", [
    "R = -5.0 if task fails",
    "R = 5.0 - L̃ - 2.0Ẽ - 1.5B̃ - 1.5Z if task succeeds",
    "L̃: latency ratio (realized / deadline), clipped to [0,2]",
    "Ẽ: online-standardized energy via EMA (α=0.01)",
    "B̃: normalized PRB cost (requested / max), clipped to [0,1]",
    "Z: fallback indicator (1 if fallback triggered, 0 otherwise)",
    "QoS-oriented: prioritizes success and timely completion"
])

# Slide 20: Network Architecture
add_content_slide(prs, "MAPPO Network Architecture", [
    "Agent embedding: 16D learnable per device (MAPPO only)",
    "Agent encoder: 2-layer MLP (29→128→128) with tanh",
    "Destination encoder: 2-layer MLP (139→128→128), shared across destinations",
    "Destination head: linear projection (128→1), masked softmax",
    "PRB head: 2-layer MLP (256→128→8), conditioned on selected destination",
    "Centralized critic: 3-layer MLP (41→256→256→1) for MAPPO, (27→256→256→1) for PPO"
])

# Slide 21: Training Algorithm
add_content_slide(prs, "Training Algorithm", [
    "PPO with clipped surrogate objective, ε=0.1",
    "Immediate-return design: G_t = r_t, Â_t = r_t - V(s_t)",
    "Entropy annealing: 0.02 → 0.002 over 40 episodes",
    "Learning rate: 3×10⁻⁴ for both actor and critic",
    "4 PPO epochs per update, minibatch size 1024",
    "Fallback transitions excluded from updates",
    "Training: 40 episodes, ~2.53 hours total"
])

# Slide 22: MAPPO Training Progress
add_image_slide(prs, "MAPPO Training Progress",
    fig_base + "mappo_training_progress_40ep.png",
    "Success rate: 57.85% (ep1) → 99.91% (ep40), stable plateau after ep13")

# Slide 23: MAPPO Evaluation Seeds
add_two_column_slide(prs, "MAPPO Three-Seed Evaluation", [
    "Seed 9001: 99.8901% success",
    "Seed 9002: 99.9390% success",
    "Seed 9003: 99.9421% success",
    "",
    "Mean: 99.9237% ± 0.0238%",
    "Avg total time: 4.5284s",
    "Energy: 227.64 Wh"
], [
    "Strong cross-seed robustness",
    "All three runs tightly clustered",
    "No large deviations",
    "Consistent generalization",
    "Stable policy behavior"
])

# Slide 24: MAPPO Action Behavior
add_image_slide(prs, "MAPPO Evaluation Overview (Seed 9001)",
    fig_base + "mappo_seed9001_run_summary.png",
    "51.73% cloud, 19.81% local, 28.46% edge; PRB: 72.76% use 1-block")

# Slide 25: MAPPO Runtime Dynamics
add_image_slide(prs, "MAPPO Runtime Resource Dynamics",
    fig_base + "mappo_eval_prb_blocks.png",
    "Moderate stable PRB reservation, balanced CPU utilization across tiers")

# Slide 26: Device Count Scalability
add_image_slide(prs, "Device Count Scalability Study",
    fig_base + "ablation_devices_train_reward.png",
    "80/160/240 devices converge; 320 devices fail due to infrastructure saturation")

# Slide 27: Device Count Results
add_two_column_slide(prs, "Device Count Evaluation Results", [
    "80 devices: 99.75% success",
    "160 devices: 99.77% success",
    "240 devices: 99.55% success",
    "320 devices: 54.06% success",
    "",
    "Graceful scaling 80→240",
    "Sharp collapse at 320"
], [
    "Delay failures:",
    "• 80: 86",
    "• 160: 154",
    "• 240: 451",
    "• 320: 61,899 (137× jump)",
    "",
    "Infrastructure saturation,",
    "not optimization failure"
])

# Slide 28: Ablation Study Summary
add_two_column_slide(prs, "Ablation Study: NoEmbedding vs FixedPRB", [
    "NoEmbedding:",
    "• Success: 99.8148%",
    "• Delay failures: 123",
    "• Energy: 227.73 Wh",
    "• Small controlled degradation",
    "• Removes 16D agent embedding",
    "• Slightly more aggressive PRB"
], [
    "FixedPRB (10 blocks):",
    "• Success: 92.3822%",
    "• Delay failures: 5127.67",
    "• Energy: 309.26 Wh",
    "• Structural failure",
    "• Disables adaptive PRB",
    "• Systematic mismatch"
])

# Slide 29: Ablation Conclusion
add_content_slide(prs, "Ablation Study Conclusions", [
    "Adaptive PRB control is the DOMINANT contributor to performance",
    "Removing PRB adaptation: 99.92% → 92.38% success, 7.5pp drop",
    "Removing agent embedding: 99.92% → 99.81% success, 0.11pp drop",
    "Agent embeddings provide smaller but consistent reliability refinement",
    "MAPPO value = joint offloading-PRB learning + agent conditioning",
    "FixedPRB shows systematic inability to manage compute-radio contention"
])

# Slide 30: Baseline Comparison Table
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Offline Baseline Comparison (Fixed Seed)"
left = Inches(0.5)
top = Inches(1.5)
width = Inches(9)
height = Inches(0.4)
textbox = slide.shapes.add_textbox(left, top, width, height)
tf = textbox.text_frame
tf.text = "MAPPO: 99.96% | PPO: 99.94% | ROUND_ROBIN: 90.39% | CLOUD: 64.64% | LOCAL: 45.36%"
tf.paragraphs[0].font.size = Pt(16)
tf.paragraphs[0].font.bold = True

add_content_slide(prs, "Baseline Comparison Highlights", [
    "MAPPO: 99.96% success, 4.56s time, 231.26 Wh, 20 delay failures",
    "PPO: 99.94% success, 4.49s time, 231.38 Wh, 36 delay failures",
    "ROUND_ROBIN: 90.39% success, 5.83s time, 298.26 Wh, 6452 failures",
    "RL policies: 9.58pp success gain, 1.28s faster, 67 Wh less energy",
    "RL bandwidth: 4.07 Mbps vs RR: 77.26 Mbps (19× reduction)",
    "Coordinated three-tier usage vs mechanical single-tier overload"
])

# Slide 31: Runtime Success Comparison
add_image_slide(prs, "Runtime Success Trajectories",
    fig_base + "mappo_runtime_success.png",
    "MAPPO/PPO: flat near 100%; ROUND_ROBIN: gradual decline to 90%")

# Slide 32: Runtime Destination Distribution
add_image_slide(prs, "Runtime Destination Distribution",
    fig_base + "mappo_runtime_destination.png",
    "MAPPO/PPO: cloud-dominant three-tier; RR: larger local fraction, flatter edge split")

# Slide 33: CPU Utilization Comparison
add_image_slide(prs, "CPU Utilization Comparison",
    fig_base + "compare_mappo_cpu_usage.png",
    "MAPPO/PPO: balanced loading; RR: edge near saturation causing queuing delays")

# Slide 34: Delay Distribution Comparison
add_image_slide(prs, "Delay Distribution Comparison",
    fig_base + "compare_mappo_delays.png",
    "MAPPO/PPO: narrow band well below deadline; RR: wider spread with tail violations")

# Slide 35: Task Failure Accumulation
add_image_slide(prs, "Task Failure Accumulation Over Time",
    fig_base + "compare_mappo_tasks_failed.png",
    "MAPPO/PPO: flat curves (20/36 failures); RR: linear accumulation (6452 failures)")

# Slide 36: MAPPO vs PPO Summary
add_two_column_slide(prs, "MAPPO vs PPO: Detailed Comparison", [
    "Three-seed means:",
    "• MAPPO: 99.9237% ± 0.0238%",
    "• PPO: 99.9436% ± 0.0137%",
    "",
    "PPO marginally stronger on:",
    "• Mean success rate",
    "• Variance",
    "• Mean total time",
    "• Mean energy"
], [
    "MAPPO advantages:",
    "• Fixed-seed: 20 vs 36 failures",
    "• Better deadline reliability",
    "• Structurally aligned",
    "• Per-device formulation",
    "",
    "MAPPO = methodologically",
    "better-matched scheduler"
])

# Slide 37: Research Question 1 Answer
add_content_slide(prs, "Answer to Research Question 1", [
    "Q1: Can on-policy RL improve QoS and system efficiency over heuristics?",
    "",
    "YES - Strong positive evidence:",
    "• MAPPO/PPO: 99.9%+ success vs 90.39% best heuristic",
    "• 9.58pp success gain, 1.28s faster, 67 Wh less energy",
    "• 19× reduction in bandwidth occupation (4.07 vs 77.26 Mbps)",
    "• Joint destination-PRB control improves both QoS and efficiency"
])

# Slide 38: Research Question 2 Answer
add_content_slide(prs, "Answer to Research Question 2", [
    "Q2: Is adaptive PRB allocation the dominant source of gain?",
    "",
    "YES - Ablation evidence is clear:",
    "• FixedPRB: 99.92% → 92.38% success (7.5pp drop)",
    "• NoEmbedding: 99.92% → 99.81% success (0.11pp drop)",
    "• Adaptive PRB control is 68× more important than agent embedding",
    "• Joint offloading-PRB learning is the key principle"
])

# Slide 39: Research Question 3 Answer
add_content_slide(prs, "Answer to Research Question 3", [
    "Q3: Does MAPPO provide meaningful advantage over PPO?",
    "",
    "QUALIFIED YES - Narrower reliability-oriented advantage:",
    "• PPO marginally stronger on aggregate three-seed means",
    "• MAPPO: fewer residual deadline failures (20 vs 36 fixed-seed)",
    "• NoEmbedding ablation: agent conditioning improves reliability",
    "• MAPPO = structurally better-matched per-device formulation",
    "• Not universally superior, but methodologically aligned"
])

# Slide 40: Key Findings Summary
add_content_slide(prs, "Key Findings Summary", [
    "On-policy RL substantially outperforms heuristic schedulers",
    "Adaptive PRB control is the dominant design factor",
    "MAPPO provides reliability-oriented advantage over PPO",
    "Graceful scaling 80→240 devices, collapse at 320 (infrastructure limit)",
    "Joint communication-computation control is essential",
    "MAPPO = methodologically aligned per-device scheduler"
])

# Slide 41: Limitations
add_content_slide(prs, "Limitations", [
    "Simulation-based evidence, gap to real 5G/6G deployment",
    "Claims tied to tested workload range and infrastructure",
    "320-device saturation shows capacity boundary, not infinite scalability",
    "Full-task offloading only, no partial offloading",
    "Limited to LOCAL-EDGE-CLOUD, no mist layer or device-to-device",
    "PPO artifacts lack per-decision trajectory exports for symmetric comparison"
])

# Slide 42: Future Work
add_content_slide(prs, "Future Work", [
    "Map capacity boundary systematically across device density and infrastructure",
    "Preserve symmetric behavioral artifacts for PPO and MAPPO",
    "Extend to partial offloading with joint task partitioning",
    "Expand architecture to include mist layer and neighboring-device execution",
    "Reduce simulation-to-reality gap with realistic traces or hardware testbed",
    "Broader robustness validation across seeds, traffic patterns, topologies"
])

# Slide 43: Conclusion
add_content_slide(prs, "Conclusion", [
    "Joint task offloading and PRB allocation formulated as Dec-POMDP",
    "MAPPO framework with CTDE achieves 99.96% success rate",
    "Adaptive PRB control is the dominant contributor to performance",
    "MAPPO provides reliability-oriented advantage, structurally aligned",
    "On-policy RL is effective within tested capacity envelope",
    "Extended PureEdgeSim provides end-to-end training/evaluation pipeline"
])

# Slide 44: Thank You
add_title_slide(prs, "Thank You", "Questions?")

# Save presentation
output_path = "EE5003_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
