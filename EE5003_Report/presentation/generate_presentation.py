from pathlib import Path
from typing import Iterable, List, Tuple

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


BASE_DIR = Path(__file__).resolve().parent.parent
FIG_DIR = BASE_DIR / "figures"
OUT_DIR = BASE_DIR / "presentation"
OUT_PATH = OUT_DIR / "EE5003_Kan_Ziyang_Presentation.pptx"


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

COLORS = {
    "navy": RGBColor(13, 43, 89),
    "red": RGBColor(189, 21, 34),
    "gold": RGBColor(196, 140, 0),
    "ink": RGBColor(36, 39, 44),
    "muted": RGBColor(102, 110, 119),
    "light_bg": RGBColor(247, 245, 242),
    "panel": RGBColor(255, 255, 255),
    "line": RGBColor(223, 226, 230),
    "soft_red": RGBColor(252, 238, 240),
    "soft_blue": RGBColor(237, 243, 252),
    "soft_gold": RGBColor(252, 248, 233),
    "success": RGBColor(16, 114, 78),
}

TITLE_FONT = "Aptos Display"
BODY_FONT = "Aptos"


def fit_image(path: Path, box_w: int, box_h: int) -> Tuple[int, int]:
    with Image.open(path) as img:
        img_w, img_h = img.size
    ratio = min(box_w / img_w, box_h / img_h)
    return int(img_w * ratio), int(img_h * ratio)


def set_run_style(run, size: int, bold: bool = False, color: RGBColor = None, font: str = BODY_FONT):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or COLORS["ink"]


def add_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["light_bg"]
    top = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.36))
    top.fill.solid()
    top.fill.fore_color.rgb = COLORS["navy"]
    top.line.fill.background()

    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, Inches(0.36), Inches(2.0), Inches(0.07))
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLORS["red"]
    accent.line.fill.background()


def add_footer(slide, number: int):
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(7.02), Inches(12.2), Inches(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["line"]
    line.line.fill.background()

    footer = slide.shapes.add_textbox(Inches(0.6), Inches(7.06), Inches(10.2), Inches(0.22))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "EE5003 MSc Project Presentation | National University of Singapore"
    set_run_style(run, 9, color=COLORS["muted"])

    num = slide.shapes.add_textbox(Inches(12.0), Inches(7.02), Inches(0.7), Inches(0.28))
    tf = num.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = str(number)
    set_run_style(run, 11, bold=True, color=COLORS["navy"])


def add_title(slide, title: str, subtitle: str = None, tag: str = None):
    if tag:
        chip = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(0.56), Inches(1.55), Inches(0.34)
        )
        chip.fill.solid()
        chip.fill.fore_color.rgb = COLORS["soft_red"]
        chip.line.color.rgb = COLORS["red"]
        chip.line.width = Pt(1)
        tf = chip.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = tag.upper()
        set_run_style(run, 10, bold=True, color=COLORS["red"])

    box = slide.shapes.add_textbox(Inches(0.65), Inches(0.92), Inches(7.2), Inches(0.7))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    set_run_style(run, 24, bold=True, color=COLORS["ink"], font=TITLE_FONT)

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.68), Inches(1.55), Inches(7.2), Inches(0.35))
        tf = sub.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = subtitle
        set_run_style(run, 11, color=COLORS["muted"])


def add_panel(slide, left, top, width, height, fill="panel", line=True):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS[fill]
    if line:
        shape.line.color.rgb = COLORS["line"]
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_bullets(slide, left, top, width, height, bullets: Iterable[str], font_size: int = 17, level0_space: int = 4):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.vertical_anchor = MSO_ANCHOR.TOP

    for idx, text in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = text
        p.level = 0
        p.space_after = Pt(level0_space)
        p.bullet = True
        for run in p.runs:
            set_run_style(run, font_size, color=COLORS["ink"])
    return box


def add_metric_box(slide, left, top, width, height, label, value, fill):
    shape = add_panel(slide, left, top, width, height, fill=fill, line=False)
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = value
    set_run_style(r1, 22, bold=True, color=COLORS["navy"], font=TITLE_FONT)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = label
    set_run_style(r2, 10, bold=True, color=COLORS["muted"])
    return shape


def add_key_message(slide, left, top, width, height, title, body, fill="soft_gold"):
    shape = add_panel(slide, left, top, width, height, fill=fill, line=False)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = title
    set_run_style(r1, 12, bold=True, color=COLORS["red"], font=TITLE_FONT)
    p1.space_after = Pt(2)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run()
    r2.text = body
    set_run_style(r2, 11, color=COLORS["ink"])


def add_image(slide, image_name: str, left, top, width, height):
    path = FIG_DIR / image_name
    new_w, new_h = fit_image(path, width, height)
    x = left + int((width - new_w) / 2)
    y = top + int((height - new_h) / 2)
    slide.shapes.add_picture(str(path), x, y, width=new_w, height=new_h)


def add_section_title_slide(prs, number, title, subtitle, logo_name="nus.jpg"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["navy"]

    red_band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, Inches(0.52), Inches(13.333), Inches(0.22))
    red_band.fill.solid()
    red_band.fill.fore_color.rgb = COLORS["red"]
    red_band.line.fill.background()

    logo = FIG_DIR / logo_name
    if logo.exists():
        slide.shapes.add_picture(str(logo), Inches(0.8), Inches(0.9), height=Inches(0.9))

    title_box = slide.shapes.add_textbox(Inches(0.85), Inches(2.05), Inches(11.2), Inches(1.0))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    set_run_style(r, 28, bold=True, color=RGBColor(255, 255, 255), font=TITLE_FONT)

    sub_box = slide.shapes.add_textbox(Inches(0.88), Inches(3.05), Inches(9.8), Inches(0.55))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = subtitle
    set_run_style(r, 15, color=RGBColor(231, 236, 245))

    footer = slide.shapes.add_textbox(Inches(0.9), Inches(6.55), Inches(8.2), Inches(0.3))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Presentation deck generated from the finalized EE5003 report"
    set_run_style(r, 11, color=RGBColor(215, 223, 236))

    add_footer(slide, number)
    return slide


def add_content_slide(prs, number, title, subtitle=None, tag=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, title, subtitle=subtitle, tag=tag)
    add_footer(slide, number)
    return slide


def main():
    OUT_DIR.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_no = 1

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["light_bg"]

    banner = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(1.15))
    banner.fill.solid()
    banner.fill.fore_color.rgb = COLORS["navy"]
    banner.line.fill.background()

    red_bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(1.15), Inches(13.333), Inches(0.12))
    red_bar.fill.solid()
    red_bar.fill.fore_color.rgb = COLORS["red"]
    red_bar.line.fill.background()

    logo = FIG_DIR / "nus.jpg"
    if logo.exists():
        slide.shapes.add_picture(str(logo), Inches(8.95), Inches(0.24), height=Inches(0.68))

    title_box = slide.shapes.add_textbox(Inches(0.78), Inches(1.72), Inches(8.8), Inches(1.7))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "On-Policy Multi-Agent Reinforcement Learning for 6G Edge-Cloud Computing"
    set_run_style(r, 28, bold=True, color=COLORS["ink"], font=TITLE_FONT)

    sub_box = slide.shapes.add_textbox(Inches(0.82), Inches(3.38), Inches(7.4), Inches(1.3))
    tf = sub_box.text_frame
    lines = [
        ("EE5003 MSc Project Report Presentation", 15, True, COLORS["red"]),
        ("Kan Ziyang | A0326643A", 13, False, COLORS["ink"]),
        ("Supervisor: Prof. Tham Chen Khong | March 2026", 12, False, COLORS["muted"]),
    ]
    for idx, (text, size, bold, color) in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = text
        set_run_style(r, size, bold=bold, color=color, font=TITLE_FONT if idx == 0 else BODY_FONT)
        p.space_after = Pt(5)

    add_panel(slide, Inches(0.82), Inches(5.15), Inches(5.75), Inches(1.1), fill="soft_blue", line=False)
    box = slide.shapes.add_textbox(Inches(1.05), Inches(5.42), Inches(5.1), Inches(0.55))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "Joint task offloading and PRB allocation in a 3-tier LOCAL-EDGE-CLOUD system"
    set_run_style(r, 14, bold=True, color=COLORS["navy"])

    add_image(slide, "extended_pureedgesim_rl_architecture.png", Inches(7.15), Inches(1.72), Inches(5.55), Inches(4.8))
    add_footer(slide, slide_no)
    slide_no += 1

    slide = add_content_slide(
        prs,
        slide_no,
        "Motivation and Problem",
        "Why joint compute-radio scheduling matters in the studied 6G edge-cloud setting",
        tag="Background",
    )
    add_panel(slide, Inches(0.62), Inches(1.95), Inches(5.35), Inches(4.72))
    add_bullets(
        slide,
        Inches(0.85),
        Inches(2.15),
        Inches(4.85),
        Inches(4.1),
        [
            "6G-oriented services require low latency, high reliability, and support for dense device populations.",
            "Pure offloading decisions are insufficient because uplink PRB scarcity directly changes transmission delay and deadline success.",
            "Devices, edge data centers, and the cloud interact through shared compute queues, wireless contention, and mobility.",
            "Static heuristics struggle because the control problem is high-dimensional, non-stationary, and partially observable.",
            "Target: jointly choose execution destination and PRB reservation to improve QoS while limiting energy and bandwidth pressure.",
        ],
        font_size=16,
    )
    add_key_message(
        slide,
        Inches(0.88),
        Inches(5.98),
        Inches(4.85),
        Inches(0.55),
        "Core challenge",
        "A destination that looks compute-efficient can still fail if the wireless reservation is badly matched to the task and link conditions.",
    )
    add_panel(slide, Inches(6.18), Inches(1.95), Inches(6.45), Inches(4.72))
    add_image(slide, "network_topology.png", Inches(6.35), Inches(2.15), Inches(6.1), Inches(4.15))
    slide_no += 1

    slide = add_content_slide(
        prs,
        slide_no,
        "Research Questions and Contributions",
        "What the report tries to answer and what was built to answer it",
        tag="Scope",
    )
    add_panel(slide, Inches(0.62), Inches(1.98), Inches(5.9), Inches(4.78))
    add_bullets(
        slide,
        Inches(0.9),
        Inches(2.2),
        Inches(5.3),
        Inches(2.3),
        [
            "RQ1: Can on-policy RL outperform heuristic schedulers when destination selection and PRB reservation are optimized jointly?",
            "RQ2: Is adaptive PRB control the dominant source of gain relative to other MAPPO design choices?",
            "RQ3: Does MAPPO offer a meaningful reliability advantage over a shared-policy PPO baseline?",
        ],
        font_size=16,
    )
    add_key_message(
        slide,
        Inches(0.9),
        Inches(5.3),
        Inches(5.15),
        Inches(1.08),
        "Contribution set",
        "Dec-POMDP formulation, CTDE MAPPO scheduler, PPO baseline, PRB-aware PureEdgeSim extensions, and a Java-Python training/evaluation pipeline.",
        fill="soft_gold",
    )
    add_panel(slide, Inches(6.72), Inches(1.98), Inches(5.88), Inches(4.78))
    add_image(slide, "mappo_ctde_framework.png", Inches(6.95), Inches(2.15), Inches(5.45), Inches(4.38))
    slide_no += 1

    slide = add_content_slide(
        prs,
        slide_no,
        "Method and Experimental Setup",
        "System structure, action space, and evaluation protocol",
        tag="Method",
    )
    add_metric_box(slide, Inches(0.68), Inches(2.0), Inches(2.1), Inches(0.82), "devices", "160", "soft_red")
    add_metric_box(slide, Inches(2.98), Inches(2.0), Inches(2.1), Inches(0.82), "edge DCs", "4", "soft_blue")
    add_metric_box(slide, Inches(5.28), Inches(2.0), Inches(2.1), Inches(0.82), "cloud nodes", "1", "soft_gold")
    add_metric_box(slide, Inches(7.58), Inches(2.0), Inches(2.1), Inches(0.82), "PRBs", "5000", "soft_red")
    add_metric_box(slide, Inches(9.88), Inches(2.0), Inches(2.1), Inches(0.82), "episode horizon", "30 min", "soft_blue")

    add_panel(slide, Inches(0.64), Inches(3.08), Inches(6.0), Inches(3.42))
    add_bullets(
        slide,
        Inches(0.9),
        Inches(3.3),
        Inches(5.45),
        Inches(2.8),
        [
            "Each task-generating device is one agent under CTDE training.",
            "Actor uses a dual-head decision: destination plus PRB bin.",
            "Observation design: 13-D agent features and 11-D per-destination features.",
            "PRB action space has 8 bins, mapped to {1, 2, 5, 10, 20, 30, 40, 50} blocks.",
            "Training: 40 episodes; evaluation: seeds 9001, 9002, 9003; comparisons include heuristics, PPO, and controlled ablations.",
        ],
        font_size=15,
    )
    add_panel(slide, Inches(6.84), Inches(3.08), Inches(5.76), Inches(3.42))
    add_image(slide, "actor-critic.png", Inches(7.02), Inches(3.3), Inches(5.4), Inches(2.95))
    slide_no += 1

    slide = add_section_title_slide(
        prs,
        slide_no,
        "Results",
        "This section is intentionally more detailed: convergence, robustness, behavior, scalability, ablation, and baseline comparison.",
    )
    slide_no += 1

    slide = add_content_slide(
        prs,
        slide_no,
        "Results 1: MAPPO Training Convergence",
        "Standard setting: 160 devices, 30-minute horizon, 40 training episodes",
        tag="Results",
    )
    add_panel(slide, Inches(0.62), Inches(1.98), Inches(7.35), Inches(4.92))
    add_image(slide, "mappo_training_progress_40ep.png", Inches(0.82), Inches(2.15), Inches(6.95), Inches(4.45))
    add_panel(slide, Inches(8.15), Inches(1.98), Inches(4.45), Inches(4.92))
    add_bullets(
        slide,
        Inches(8.38),
        Inches(2.18),
        Inches(3.95),
        Inches(3.25),
        [
            "Task success rises from 57.85% in episode 1 to 99.91% in episode 40.",
            "MAPPO exceeds 95% success by episode 6 and 99% by episode 13, indicating fast learning under the given action space.",
            "Average total time drops from 5.0108 s to 4.6622 s while delay failures fall from 28,193 to 52.",
            "Late-stage training is stable: last-10-episode success averages 99.80% with 0.131 std.",
            "Training cost is practical for this simulator scale: 227.74 s per episode, about 2.53 hours total.",
        ],
        font_size=14,
    )
    add_key_message(
        slide,
        Inches(8.38),
        Inches(5.82),
        Inches(3.9),
        Inches(0.72),
        "Interpretation",
        "The policy does not only learn to finish more tasks; it simultaneously reduces latency pressure and stabilizes the operating regime.",
    )
    slide_no += 1

    slide = add_content_slide(
        prs,
        slide_no,
        "Results 2: Robustness and Behavioral Pattern",
        "Three-seed evaluation plus representative seed-9001 behavior",
        tag="Results",
    )
    add_metric_box(slide, Inches(0.72), Inches(2.02), Inches(2.5), Inches(0.9), "mean success", "99.9237%", "soft_red")
    add_metric_box(slide, Inches(3.45), Inches(2.02), Inches(2.5), Inches(0.9), "success std", "0.0238", "soft_blue")
    add_metric_box(slide, Inches(6.18), Inches(2.02), Inches(2.5), Inches(0.9), "avg total time", "4.5284 s", "soft_gold")
    add_metric_box(slide, Inches(8.91), Inches(2.02), Inches(2.5), Inches(0.9), "mean energy", "227.636 Wh", "soft_red")
    add_panel(slide, Inches(0.64), Inches(3.18), Inches(7.04), Inches(3.4))
    add_bullets(
        slide,
        Inches(0.92),
        Inches(3.4),
        Inches(6.5),
        Inches(2.9),
        [
            "All three evaluation seeds remain tightly clustered: 99.8901%, 99.9390%, and 99.9421% success.",
            "Representative seed 9001 shows a balanced destination policy: 51.73% cloud, 19.81% local, and 28.46% edge.",
            "PRB behavior is selective rather than aggressive: 1-block is used for 72.76% of offloaded tasks, while 10/20 blocks are reserved only when needed.",
            "Runtime reward stays around 4.4-4.5 after initialization, which indicates stable task-level decisions across the full 30-minute horizon.",
        ],
        font_size=14,
    )
    add_panel(slide, Inches(7.9), Inches(3.18), Inches(4.7), Inches(3.4))
    add_image(slide, "mappo_seed9001_run_summary.png", Inches(8.08), Inches(3.34), Inches(4.35), Inches(3.08))
    slide_no += 1

    slide = add_content_slide(
        prs,
        slide_no,
        "Results 3: Runtime Resource Dynamics",
        "Why the converged policy stays reliable in operation",
        tag="Results",
    )
    add_panel(slide, Inches(0.62), Inches(1.98), Inches(5.95), Inches(4.72))
    add_image(slide, "mappo_eval_prb_blocks.png", Inches(0.8), Inches(2.15), Inches(5.55), Inches(2.0))
    add_image(slide, "mappo_eval_cpu_usage.png", Inches(0.8), Inches(4.22), Inches(5.55), Inches(2.2))
    add_panel(slide, Inches(6.82), Inches(1.98), Inches(5.78), Inches(4.72))
    add_bullets(
        slide,
        Inches(7.08),
        Inches(2.18),
        Inches(5.2),
        Inches(3.1),
        [
            "PRB allocation remains moderate and smooth, matching the action statistics where small reservations dominate.",
            "CPU load is spread across cloud and edge tiers while local utilization remains comparatively low.",
            "The delay distribution has a thin tail and the failure curve stays almost flat, so the near-perfect success rate is not hiding bursts of instability.",
            "Operationally, MAPPO is coordinating compute placement and radio reservation together instead of solving them as disconnected subproblems.",
        ],
        font_size=14,
    )
    add_key_message(
        slide,
        Inches(7.08),
        Inches(5.82),
        Inches(5.0),
        Inches(0.72),
        "Takeaway",
        "The learned policy uses cloud and edge capacity without oversaturating the wireless layer or pushing too much traffic back to local devices.",
        fill="soft_blue",
    )
    slide_no += 1

    slide = add_content_slide(
        prs,
        slide_no,
        "Results 4: Scalability Across Device Counts",
        "Independent training and evaluation at 80, 160, 240, and 320 devices",
        tag="Results",
    )
    add_panel(slide, Inches(0.62), Inches(1.98), Inches(6.15), Inches(4.72))
    add_image(slide, "ablation_devices_success_rate.png", Inches(0.82), Inches(2.16), Inches(5.85), Inches(2.0))
    add_image(slide, "ablation_devices_eval_reward.png", Inches(0.82), Inches(4.2), Inches(5.85), Inches(2.15))
    add_panel(slide, Inches(7.0), Inches(1.98), Inches(5.6), Inches(4.72))
    add_bullets(
        slide,
        Inches(7.25),
        Inches(2.16),
        Inches(5.05),
        Inches(3.38),
        [
            "MAPPO scales well from 80 to 240 devices: success stays at 99.75%, 99.77%, and 99.55%, respectively.",
            "At 320 devices, success collapses to 54.06% and delay failures jump to 61,899, so the workload crosses the infrastructure capacity boundary.",
            "The policy adapts before collapse: cloud share drops from about 60% at 80 devices to 11.7% at 320 devices, while local share rises to 50.7%.",
            "This pattern indicates infrastructure saturation rather than ordinary optimization noise; the system is running out of compute-radio capacity, not simply failing to train.",
        ],
        font_size=14,
    )
    add_key_message(
        slide,
        Inches(7.25),
        Inches(5.92),
        Inches(5.02),
        Inches(0.62),
        "Capacity-envelope conclusion",
        "Within the tested four-edge-plus-one-cloud infrastructure, MAPPO is reliable up to roughly 240 devices.",
        fill="soft_gold",
    )
    slide_no += 1

    slide = add_content_slide(
        prs,
        slide_no,
        "Results 5: Ablation Study",
        "Is MAPPO winning because of multi-agent conditioning or because of adaptive PRB control?",
        tag="Results",
    )
    add_metric_box(slide, Inches(0.72), Inches(2.0), Inches(3.5), Inches(0.88), "MAPPO mean success", "99.9237%", "soft_red")
    add_metric_box(slide, Inches(4.5), Inches(2.0), Inches(3.5), Inches(0.88), "NoEmbedding success", "99.8148%", "soft_blue")
    add_metric_box(slide, Inches(8.28), Inches(2.0), Inches(4.0), Inches(0.88), "FixedPRB success", "92.3822%", "soft_gold")
    add_panel(slide, Inches(0.62), Inches(3.15), Inches(6.0), Inches(3.55))
    add_bullets(
        slide,
        Inches(0.9),
        Inches(3.4),
        Inches(5.45),
        Inches(2.9),
        [
            "NoEmbedding removes the 16-D agent embedding but keeps the same policy structure.",
            "Its degradation is small but real: delay failures rise from 48 to 123 and the policy shifts a bit more mass toward 20-block PRB actions.",
            "This means explicit agent identity mainly sharpens reliability under shared multi-device pressure.",
        ],
        font_size=14,
    )
    add_panel(slide, Inches(6.86), Inches(3.15), Inches(5.74), Inches(3.55))
    add_bullets(
        slide,
        Inches(7.12),
        Inches(3.4),
        Inches(5.16),
        Inches(2.9),
        [
            "FixedPRB clamps every offloaded task to 10 blocks, removing adaptive radio control.",
            "That causes structural failure: energy rises to 309.2645 Wh and delay failures explode to 5127.67.",
            "Conclusion: adaptive PRB control is the dominant source of performance gain; embeddings are a secondary but useful refinement.",
        ],
        font_size=14,
    )
    add_image(slide, "fixedprb_run_summary.png", Inches(8.55), Inches(4.4), Inches(3.55), Inches(1.95))
    slide_no += 1

    slide = add_content_slide(
        prs,
        slide_no,
        "Results 6: RL vs Heuristic Baselines",
        "Fixed-seed offline comparison under the same realized 160-device workload",
        tag="Results",
    )
    add_panel(slide, Inches(0.62), Inches(1.98), Inches(7.0), Inches(4.72))
    add_image(slide, "baseline_comparison_metrics.png", Inches(0.82), Inches(2.18), Inches(6.6), Inches(4.3))
    add_panel(slide, Inches(7.85), Inches(1.98), Inches(4.75), Inches(4.72))
    add_bullets(
        slide,
        Inches(8.12),
        Inches(2.18),
        Inches(4.18),
        Inches(3.1),
        [
            "MAPPO reaches 99.96% success and PPO 99.94%, while the strongest heuristic, ROUND_ROBIN, reaches 90.39%.",
            "MAPPO beats ROUND_ROBIN by 9.58 percentage points while also cutting average total time by 1.2750 s.",
            "Energy is lower for RL as well: about 231 Wh for MAPPO/PPO versus 298.26 Wh for ROUND_ROBIN.",
            "Bandwidth usage reveals another systems gain: MAPPO and PPO use only 4.07 Mbps on average, far below heuristic schedulers.",
            "The improvement comes from balanced tier usage, not from simply pushing more tasks to the cloud or spending more energy.",
        ],
        font_size=14,
    )
    slide_no += 1

    slide = add_content_slide(
        prs,
        slide_no,
        "Results 7: MAPPO vs PPO",
        "How to interpret the two RL policies without overstating the claim",
        tag="Results",
    )
    add_panel(slide, Inches(0.62), Inches(1.98), Inches(6.1), Inches(4.72))
    add_image(slide, "mappo_ppo_eval_summary.png", Inches(0.82), Inches(2.18), Inches(5.75), Inches(4.2))
    add_panel(slide, Inches(6.98), Inches(1.98), Inches(5.62), Inches(4.72))
    add_bullets(
        slide,
        Inches(7.24),
        Inches(2.18),
        Inches(5.08),
        Inches(3.2),
        [
            "PPO is marginally stronger on three-seed mean metrics: 99.9436% success, 4.4665 s average time, and 227.5745 Wh energy.",
            "However, the fixed-seed offline run still favors MAPPO on residual failures: 20 for MAPPO versus 36 for PPO.",
            "The NoEmbedding ablation supports the same interpretation: agent-specific conditioning improves deadline-sensitive reliability.",
            "Therefore, MAPPO should be positioned as the structurally better-matched per-device formulation, while PPO remains a very strong learning baseline.",
        ],
        font_size=14,
    )
    add_key_message(
        slide,
        Inches(7.24),
        Inches(5.88),
        Inches(5.0),
        Inches(0.64),
        "Balanced claim",
        "The report supports MAPPO as a reliability-oriented and methodologically aligned scheduler, not as an unconditional winner on every aggregate metric.",
        fill="soft_red",
    )
    slide_no += 1

    slide = add_content_slide(
        prs,
        slide_no,
        "Conclusion",
        "What the presentation should leave the audience with",
        tag="Close",
    )
    add_panel(slide, Inches(0.62), Inches(1.98), Inches(12.0), Inches(4.72))
    add_bullets(
        slide,
        Inches(0.9),
        Inches(2.2),
        Inches(11.4),
        Inches(2.85),
        [
            "On-policy reinforcement learning is effective for joint offloading and PRB allocation in the studied LOCAL-EDGE-CLOUD simulator setting.",
            "The main reason is adaptive compute-radio coordination: removing adaptive PRB control causes the largest drop in success and the largest rise in energy and delay failures.",
            "MAPPO is best interpreted as a reliability-oriented, per-device formulation that matches the system structure well, while PPO serves as a strong aggregate-performance baseline.",
            "The learned policy remains highly effective from 80 to 240 devices, but performance collapses at 320 devices because the infrastructure is saturated.",
        ],
        font_size=16,
    )
    add_key_message(
        slide,
        Inches(0.9),
        Inches(5.55),
        Inches(3.55),
        Inches(0.75),
        "Future direction 1",
        "Map the capacity boundary by jointly scaling device density, compute capacity, and PRB budgets.",
        fill="soft_blue",
    )
    add_key_message(
        slide,
        Inches(4.72),
        Inches(5.55),
        Inches(3.55),
        Inches(0.75),
        "Future direction 2",
        "Preserve symmetric PPO/MAPPO trajectory artifacts and expand robustness tests across more seeds and traffic patterns.",
        fill="soft_gold",
    )
    add_key_message(
        slide,
        Inches(8.54),
        Inches(5.55),
        Inches(3.55),
        Inches(0.75),
        "Future direction 3",
        "Extend from full-task offloading to partial offloading and broader continuum architectures.",
        fill="soft_red",
    )
    slide_no += 1

    prs.save(str(OUT_PATH))
    print(f"Saved presentation to: {OUT_PATH}")


if __name__ == "__main__":
    main()
